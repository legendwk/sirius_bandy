from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
import codecs

prs=Presentation()

# 
lyt=prs.slide_layouts[0] # choosing a slide layout


slide=prs.slides.add_slide(lyt) # adding a slide
title=slide.shapes.title # assigning a title
subtitle=slide.placeholders[1] # placeholder for subtitle
title.text="Hey,This is a Slide! How exciting!" # title
subtitle.text="Really?" # subtitle
prs.save("slide1.pptx") # saving file